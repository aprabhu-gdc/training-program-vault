import argparse
import sys
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".docx", ".pdf", ".pptx", ".xlsx", ".xlsm"}


if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")


def extract_docx(path: Path) -> str:
    doc = Document(path)
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                parts.append(" | ".join(values))
    return "\n".join(parts)


def extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text.strip())
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    parts = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts = [f"Slide {index}"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            text = text.strip()
            if text:
                slide_parts.append(text)
        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))
    return "\n\n".join(parts)


def extract_xlsx(path: Path, max_rows: int = 60, max_cols: int = 12) -> str:
    workbook = load_workbook(filename=str(path), data_only=True, read_only=True)
    parts = []
    for worksheet in workbook.worksheets:
        parts.append(f"Worksheet: {worksheet.title}")
        row_count = 0
        for row in worksheet.iter_rows(min_row=1, max_row=max_rows, max_col=max_cols, values_only=True):
            values = []
            for value in row:
                if value is None:
                    continue
                text = str(value).strip()
                if text:
                    values.append(text)
            if values:
                parts.append(" | ".join(values))
                row_count += 1
        if row_count == 0:
            parts.append("[no non-empty rows in sampled range]")
    workbook.close()
    return "\n".join(parts)


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix in {".xlsx", ".xlsm"}:
        return extract_xlsx(path)
    raise ValueError(f"Unsupported extension: {suffix}")


def iter_files(paths: list[Path], recursive: bool) -> list[Path]:
    files: list[Path] = []
    for path in paths:
        if path.is_dir():
            pattern = "**/*" if recursive else "*"
            for file_path in sorted(path.glob(pattern)):
                if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                    files.append(file_path)
        elif path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            files.append(path)
    return files


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract text from Office and PDF files.")
    parser.add_argument("paths", nargs="+", help="File or folder paths")
    parser.add_argument("--recursive", action="store_true", help="Walk folders recursively")
    parser.add_argument("--chars", type=int, default=6000, help="Maximum characters per file")
    args = parser.parse_args()

    files = iter_files([Path(path) for path in args.paths], args.recursive)
    if not files:
        print("No supported files found.", file=sys.stderr)
        return 1

    for file_path in files:
        print(f"=== {file_path} ===")
        try:
            text = extract_text(file_path)
        except Exception as exc:
            print(f"[extract-error] {exc}")
        else:
            text = (text or "").strip()
            if not text:
                print("[no extracted text]")
            else:
                print(text[: args.chars])
        print()
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
